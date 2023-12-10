from pptx import Presentation
from pptx import *
import os
import openpyxl
import xlsxwriter




def open_output(prs_name):
    try:
        wb = openpyxl.load_workbook(f"{prs_name}_insides.xlsx")
    except Exception as error:
        print(error)
        f = xlsxwriter.Workbook(f"{prs_name}_insides.xlsx")
        f.close()
        wb = openpyxl.load_workbook(f"{prs_name}_insides.xlsx")
    return wb

def separate(prs):
    file = {}
    for slide in prs.slides:
        #print(f"Slide {prs.slides.index(slide)}")
        file[f"Slide {prs.slides.index(slide)}"] = {}
        shape_nb = 1
        for shape in slide.shapes:
            #print("  ", shape.__class__.__name__)
            file[f"Slide {prs.slides.index(slide)}"][f"{shape.name}"] = ""
            if shape.__class__.__name__ == "Picture":
                pass
            else:
                try:
                    #print("  Text:", shape.text)
                    file[f"Slide {prs.slides.index(slide)}"][f"{shape.name}"] = shape.text.replace("\n", " ")
                except:
                    pass
            shape_nb += 1
    
    return file

