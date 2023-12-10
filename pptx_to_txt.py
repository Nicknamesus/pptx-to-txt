from pptx import Presentation
from pptx import *
import os
import openpyxl
import re
from openpyxl.cell.cell import ILLEGAL_CHARACTERS_RE
from openpyxl.styles import Alignment
from util import *
import math

# load a presentation
prs_url = input("Input the full path to the file without the: ")
prs_url = prs_url.replace('"', '')
#prs_url= r"C:\Users\nikit\Documents\pptx_to_txt\pptx_to_txt\PRODUCTX-VersionY.pptx"
try:
    prs_folder = prs_url[0: prs_url.rfind('\\')]
    os.chdir(prs_folder)
    prs_name = prs_url[prs_url.rfind('\\')+1: prs_url.rfind(".")]
    print(prs_name)
    file_name = f"{prs_name}_insides.xlsx"
    prs = Presentation(prs_url)
except:
    print("There seems to be something wrong with the provided path. Please try again.")
    exit()

'''
print(prs)
print(prs.slides)
for slide in prs.slides:
    print("Slide", prs.slides.index(slide))
    for shape in slide.shapes:
        print("  ", shape.__class__.__name__)
        try:
            print("  Text:", shape.text)
        except:
            pass
        #print(shape.element)
'''

wb = open_output(prs_name)

if __name__ == "__main__":
    sheet = wb.active
    file_contents = separate(prs)
    
    row = 1
    for slide in file_contents:
        #putting slides
        sheet.cell(row=row, column = 1).value = slide
        row += 1
        for shape in file_contents[slide]:
            #putting slide element names
            sheet.cell(row=row, column = 2).value = shape
            #putting the text of those elements
            sheet.cell(row=row, column = 3).value = ILLEGAL_CHARACTERS_RE.sub(r'',file_contents[slide][shape])


            if len(file_contents[slide][shape]) > 200:
                wb.worksheets[0].row_dimensions[row].height = 15*(math.ceil((len(file_contents[slide][shape])/200)))
            sheet[f"C{row}"].alignment = Alignment(wrap_text=True)
            row += 1

    
    wb.worksheets[0].column_dimensions["B"].width = 30
    wb.worksheets[0].column_dimensions["C"].width = 200
    wb.save(file_name)
    